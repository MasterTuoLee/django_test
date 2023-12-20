"""
爬取http://www.cssmoban.com/ppt/上的所有ppt图片并加工成完整的ppt（这里是因为没看到网页直接有内容下载，才有使用pptx来完成图片转ppt）
"""
import os
import requests
from pptx import Presentation
from pptx.util import Inches
from pyquery import PyQuery

# 根网址
base_url = "http://www.cssmoban.com/ppt/"

# 设置异常捕获
try:
    for index in range(1, 5):
        response = requests.get(url=f"{base_url}{index}.html")
        if response.status_code == 200:
            doc_response = PyQuery(response.text)

            #
            img_lists = doc_response(".info-list")

            # 创建文件夹，指定下载路径
            save_folder = f'ppt/ppt_{index}'
            os.makedirs(save_folder, exist_ok=True)

            # 制作ppt列表
            img_files = []
            # 遍历获取每一个图像地址
            for e in img_lists:
                # 获取图像下载地址
                img_url = PyQuery(e)("img").attr("src")
                # 指定下载名
                img_name = img_url.split('/')[-1]
                # 拼接下载路径
                img_path = os.path.join(save_folder, img_name)
                img_files.append(img_path)
                # 开始请求下载
                ppt_response = requests.get(url=img_url)
                if ppt_response.status_code == 200:
                    with open(img_path, 'wb') as img_file:
                        img_file.write(ppt_response.content)
                        print(f"下载ppt{index}中的{img_name}完成")
                else:
                    print("下载失败")
            print(f"下载ppt{index}完成")

            # 创建一个ppt对象
            ppt = Presentation()

            # 获取ppt图片列表
            for img_file in img_files:
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])

                left = Inches(0)
                top = Inches(0)
                width = ppt.slide_width
                height = ppt.slide_height

                img = slide.shapes.add_picture(img_file, left, top, width, height)

            ppt.save(f"{index}.pptx")
            print(f"ppt{index}制作完成")


        else:
            print("请求失败")
except Exception as e:
    print(f"出异常了！！！异常类型为：{e}")
