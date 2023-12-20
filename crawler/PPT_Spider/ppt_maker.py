from pptx import Presentation
from pptx.util import Inches

# 创建一个ppt文件
ppt = Presentation()
# img_folder = "ppt/ppt_1"

# 图片名称列表
# img_files = [f for f in os.listdir(img_folder) if f.endswith(".jpg")]

img_files = ["preview_1.jpg",
             "preview_2.jpg",
             "preview_3.jpg",
             "preview_4.jpg",
             "preview_5.jpg", ]

for img_file in img_files:
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])

    left = Inches(0)
    top = Inches(0)
    width = ppt.slide_width
    height = ppt.slide_height

    img = slide.shapes.add_picture(img_file, left, top, width, height)

ppt.save("1.pptx")
